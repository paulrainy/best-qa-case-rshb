import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

def set_text_with_styling(shape, new_text_paragraphs, is_title=False):
    if not shape.has_text_frame:
        return
    
    tf = shape.text_frame
    tf.word_wrap = True
    
    # Capture template styling
    template_font_name = "Proxima Nova"
    template_font_size = Pt(45) if is_title else Pt(22)
    # Titles are NEON GREEN, Body is WHITE/LIGHT BLUE
    target_color = RGBColor(0x88, 0xDF, 0x59) if is_title else RGBColor(0xD8, 0xDF, 0xF6)
    
    if len(tf.paragraphs) > 0 and len(tf.paragraphs[0].runs) > 0:
        ref_run = tf.paragraphs[0].runs[0]
        if ref_run.font.name: template_font_name = ref_run.font.name
        if not is_title and ref_run.font.size: 
             template_font_size = ref_run.font.size

    tf.clear()
    
    for i, p_text in enumerate(new_text_paragraphs):
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(12) if i > 0 else Pt(0)
        
        run = p.add_run()
        run.text = p_text
        run.font.name = template_font_name
        run.font.size = template_font_size
        run.font.color.rgb = target_color

def generate_showcase():
    template_path = 'docs/template.pptx'
    output_path = 'docs/showcase.pptx'
    
    if not os.path.exists(template_path):
        print(f"Template not found: {template_path}")
        return

    prs = Presentation(template_path)
    
    # Streamlined Content for 7 Slides (0-6)
    content = {
        0: {
            "title": "Кредитный Конвейер",
            "subtitle": "Стратегия и автоматизация тестирования\nЛучший Тестировщик 2026",
            "footer": "Стек: Python + Docker + Allure"
        },
        1: {
            "title": "Решаемая проблема",
            "body": [
                "Тестирование «Черного ящика» (SP_GET_CREDIT_DECISION)",
                "Риск ошибок в критической бизнес-логике банка",
                "Сложные пересекающиеся правила (Возраст vs Доход)",
                "Непротиворечивость и полнота требований (ТЗ)"
            ]
        },
        2: {
            "title": "Стратегия тестирования",
            "body": [
                "DDT: Отделение логики тестов от данных (CSV)",
                "Техники: Decision Table + BVA + Equivalence Partitioning",
                "Traceability: Прямая связь отчетов с пунктами ТЗ",
                "100% покрытие всех бизнес-правил и исключений"
            ]
        },
        3: {
            "title": "Критический анализ ТЗ",
            "body": [
                "Бизнес-конфликт: «BAD history» vs «Исключение по возрасту»",
                "Отсутствие валидации: Не описано поведение при NULL / Invalid",
                "Граничные условия: Уточнение строгости неравенств (> vs >=)",
                "Риск: Неопределенность вердикта при конфликте правил"
            ]
        },
        4: {
            "title": "Техническая реализация",
            "body": [
                "Инфраструктура: PostgreSQL в Docker (изолированная среда)",
                "Тест-движок: Pytest + Psycopg2 + Pandas",
                "Отчетность: Allure Report с детализацией шагов БД",
                "Масштабируемость: Добавление кейсов без правки кода"
            ]
        },
        5: {
            "title": "Результаты и бизнес-ценность",
            "body": [
                "Автоматизация «под ключ»: от запуска БД до отчета",
                "Стабильность: Гарантия корректности скоринга перед продом",
                "Экономия: Быстрый регресс при изменении хранимых процедур",
                "Качество: Обнаружение 3-х «серых зон» в ТЗ до разработки"
            ]
        },
        6: {
            "title": "СПАСИБО ЗА ВНИМАНИЕ!",
            "subtitle": "Готов ответить на ваши вопросы",
            "footer": "Автор: [Ваше ФИО]"
        }
    }

    # Delete slides beyond 7 (from index 7 onwards)
    # Slides collection is not directly deletable easily in python-pptx
    # Better to just use only the slides we need if we are copying, 
    # but here we are using a template with 10 slides.
    # We will clear all text frames in slides 7-9 to make them effectively "Empty"
    # and we could potentially use a script to delete them if needed.

    print("Generating streamlined showcase.pptx (7 slides)...")
    for i, slide in enumerate(prs.slides):
        if i not in content:
            # Clear all text in unused slides
            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape.text_frame.clear()
            continue
        
        data = content[i]
        title_shape = None
        subtitle_shape = None
        body_shapes = []
        footer_shape = None
        
        # Identify shapes by keywords from template
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            text = "".join([p.text for p in shape.text_frame.paragraphs]).strip().lower()
            if not text: continue
            
            if i == 0:
                if "efr down" in text: title_shape = shape
                elif "интеллектуальная" in text: subtitle_shape = shape
                elif "трек" in text: footer_shape = shape
            elif i == 6: # New final slide (was 9) - but wait, the template has indices 0-9.
                # If we want Slide 6 to be the final, we should use the styling of slide 9.
                # Actually, I'll just move the content to Slide 6 and leave 7,8,9 empty.
                pass 
            
            # General keyword matching for titles
            if any(kw in text for kw in ["решаемая", "зачем этот", "концепция", "архитектурная", "оценка", "текущее", "потенциальная", "благодарим"]):
                title_shape = shape
            elif len(text) > 15 and shape != title_shape:
                body_shapes.append(shape)

        # Apply content
        if title_shape:
            set_text_with_styling(title_shape, [data["title"]], is_title=True)
        if "subtitle" in data:
            # For Slide 0 or Final, find the subtitle shape
            for shape in slide.shapes:
                if shape.has_text_frame and shape != title_shape:
                    txt = "".join([p.text for p in shape.text_frame.paragraphs]).lower()
                    if i == 0 and "интеллектуальная" in txt: subtitle_shape = shape
                    elif i == 6 and "готов ответить" in txt: subtitle_shape = shape
            if subtitle_shape:
                set_text_with_styling(subtitle_shape, [data["subtitle"]])
                
        if footer_shape and "footer" in data:
            set_text_with_styling(footer_shape, [data["footer"]])
            
        if "body" in data and body_shapes:
            main_body = body_shapes[0]
            set_text_with_styling(main_body, data["body"])
            for other in body_shapes[1:]:
                other.text_frame.clear()

        # Cleanup Slide 6 specifically if it's the final (we use the template's Slide 9 look for Slide 6?)
        # No, better to keep the indices as is for now but only fill up to 6.
        # Actually Slide 6 in template is "Текущее состояние".
        # Let's map Slide 6 content (Thank you) to Slide 9 but then we have a gap.
        # User asked: "зачем нам слайды 8 и 9?".
        # This implies they want Slide 7 to be the last one.
        
    # Final fix for slide deletion: we will create a NEW presentation and copy slides
    # but that's complex with styles. Simpler: tell user we cleared 7,8,9.
    # OR: Actually delete them.
    
    # Python-pptx doesn't have a simple slide.delete().
    # Here is a workaround for slide deletion:
    def delete_slide(presentation, index):
        xml_slides = presentation.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[index])

    # Delete slides 7, 8, 9 (in reverse order to keep indices stable)
    if len(prs.slides) > 7:
        for i in range(len(prs.slides) - 1, 6, -1):
            delete_slide(prs, i)

    prs.save(output_path)
    print(f"Successfully saved to {output_path} (Total slides: {len(prs.slides)})")

if __name__ == "__main__":
    generate_showcase()
